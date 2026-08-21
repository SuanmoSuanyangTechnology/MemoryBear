"""Bounded synchronous preview parsing for interface requests."""

from __future__ import annotations

import csv
import io
import json
import re
from collections.abc import Iterable
from io import BytesIO
from pathlib import Path
from typing import Any

import tiktoken

from ...errors import KnowledgeError
from ..models.chunk import ChildDocumentChunk, DocumentChunk

_ENCODER = tiktoken.get_encoding("cl100k_base")
_HTML_TAG = re.compile(r"<[^>]+>")


def _decode_text(binary: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return binary.decode(encoding)
        except UnicodeDecodeError:
            continue
    return binary.decode("utf-8", errors="replace")


def _flatten_json(value: Any) -> Iterable[str]:
    if isinstance(value, dict):
        for key, item in value.items():
            if isinstance(item, (dict, list)):
                yield from _flatten_json(item)
            elif item is not None:
                yield f"{key}: {item}"
    elif isinstance(value, list):
        for item in value:
            yield from _flatten_json(item)
    elif value is not None:
        yield str(value)


def _extract_text(filename: str, binary: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(binary))
        return "\n\n".join((page.extract_text() or "") for page in reader.pages[:5])
    if suffix == ".docx":
        from docx import Document

        document = Document(BytesIO(binary))
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        for table in document.tables:
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells if cell.text))
        return "\n".join(parts)
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(BytesIO(binary))
        return "\n\n".join(
            "\n".join(
                shape.text
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text
            )
            for slide in presentation.slides[:5]
        )
    if suffix in {".xlsx", ".xls"}:
        return _extract_spreadsheet(binary, suffix)
    if suffix == ".json":
        try:
            return "\n".join(_flatten_json(json.loads(_decode_text(binary))))
        except json.JSONDecodeError as exc:
            raise KnowledgeError.from_code(
                "KB_VALIDATION_ERROR",
                "Invalid JSON document",
            ) from exc
    if suffix in {".html", ".htm"}:
        from bs4 import BeautifulSoup

        return BeautifulSoup(_decode_text(binary), "html.parser").get_text("\n")
    if suffix == ".csv":
        rows = csv.reader(io.StringIO(_decode_text(binary)))
        return "\n".join("\t".join(cell for cell in row if cell) for row in rows)
    if suffix in {
        ".txt",
        ".md",
        ".markdown",
        ".log",
        ".xml",
        ".yaml",
        ".yml",
        ".py",
        ".js",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".php",
        ".go",
        ".ts",
        ".sh",
        ".cs",
        ".kt",
        ".sql",
        ".jsonl",
        ".ldjson",
    }:
        return _decode_text(binary)
    raise KnowledgeError.from_code(
        "KB_VALIDATION_ERROR",
        f"Preview parser is unavailable for file type: {suffix or 'unknown'}",
    )


def _extract_spreadsheet(binary: bytes, suffix: str) -> str:
    rows: list[str] = []
    if suffix == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(BytesIO(binary), read_only=True, data_only=True)
        try:
            for sheet in workbook.worksheets:
                rows.append(sheet.title)
                for row in sheet.iter_rows(values_only=True):
                    values = [str(value) for value in row if value not in (None, "")]
                    if values:
                        rows.append("\t".join(values))
        finally:
            workbook.close()
        return "\n".join(rows)

    import xlrd

    workbook = xlrd.open_workbook(file_contents=binary)
    for sheet in workbook.sheets():
        rows.append(sheet.name)
        for row_index in range(sheet.nrows):
            values = [str(value) for value in sheet.row_values(row_index) if value != ""]
            if values:
                rows.append("\t".join(values))
    return "\n".join(rows)


def _split_tokens(text: str, limit: int) -> list[str]:
    tokens = _ENCODER.encode(text)
    chunks = []
    for start in range(0, len(tokens), max(limit, 1)):
        chunks.append(_ENCODER.decode(tokens[start : start + max(limit, 1)]).strip())
    return [chunk for chunk in chunks if chunk]


def _split_normal(text: str, parser_config: dict[str, Any]) -> list[str]:
    delimiter = parser_config.get("delimiter")
    limit = max(int(parser_config.get("chunk_token_num") or 128), 1)
    parts = text.split(delimiter) if delimiter else [text]
    result: list[str] = []
    for part in parts:
        stripped = part.strip()
        if not stripped:
            continue
        if len(_ENCODER.encode(stripped)) <= limit:
            result.append(stripped)
        else:
            result.extend(_split_tokens(stripped, limit))
    return result


def _parent_child_preview(
    text: str,
    parser_config: dict[str, Any],
) -> list[DocumentChunk]:
    parent_delimiter = parser_config.get("parent_chunk_delimiter") or "\n\n"
    child_limit = max(int(parser_config.get("chunk_token_num") or 128), 1)
    parents = [part.strip() for part in text.split(parent_delimiter) if part.strip()]
    result = []
    child_sort_id = 0
    for parent_sort_id, parent in enumerate(parents):
        children = []
        for child in _split_tokens(parent, child_limit):
            children.append(
                ChildDocumentChunk(
                    page_content=child,
                    metadata={"chunk_type": "child", "sort_id": child_sort_id},
                )
            )
            child_sort_id += 1
        if not children:
            continue
        result.append(
            DocumentChunk(
                page_content=parent,
                metadata={"chunk_type": "parent", "sort_id": parent_sort_id},
                children=children,
            )
        )
    return result


def preview_binary(
    filename: str,
    binary: bytes,
    parser_config: dict[str, Any] | None,
) -> list[DocumentChunk]:
    """Parse one in-memory file without importing task orchestration."""

    config = dict(parser_config or {})
    text = _extract_text(filename, binary)
    if config.get("parent_child_mode") or config.get("chunk_mode") == "parent_child":
        return _parent_child_preview(text, config)
    return [
        DocumentChunk(
            page_content=content,
            metadata={"chunk_type": "chunk", "sort_id": index},
            children=[],
        )
        for index, content in enumerate(_split_normal(text, config))
    ]


__all__ = ["preview_binary"]
